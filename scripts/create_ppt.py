#!/usr/bin/env python3
"""
Create an academic PowerPoint presentation for the RVCE Metro Grid Digital Twin project.

Generates a complete presentation with all research plots, RVCE branding,
team details, and academic formatting. Output: docs/Metro_Grid_Digital_Twin.pptx

Usage:
    python scripts/create_ppt.py
"""

from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
RVCE_BLUE = RGBColor(0x00, 0x3D, 0x7A)
RVCE_GOLD = RGBColor(0xD4, 0xA0, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GREY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_GREEN = RGBColor(0x05, 0x96, 0x69)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)

OUTPUT_DIR = Path(__file__).resolve().parents[1] / "docs"
IMAGES_DIR = OUTPUT_DIR / "images"
OUTPUT_PATH = OUTPUT_DIR / "Metro_Grid_Digital_Twin.pptx"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

TEAM = [
    ("Varshini C. B.", "1RV22EE035"),
    ("Vedant", "1RV22EE036"),
    ("Sethu S.", "1RV22EE030"),
    ("Aravind Kumar N.", "1RV22EE007"),
]
GUIDE = "Dr. Manjunatha C."


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            a = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            a.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=12,
                       color=DARK, line_spacing=1.5):
    """Add text box with multiple formatted lines (bullets)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}" if not line.startswith("•") and not line.startswith(" ") else line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.4)
        p.level = 0
    return txBox


def add_header_bar(slide):
    """Add a top header bar with RVCE branding."""
    # Top bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9), RVCE_BLUE)
    # Gold accent line
    add_shape(slide, Inches(0), Inches(0.9), SLIDE_WIDTH, Inches(0.05), RVCE_GOLD)
    # Title text in header
    add_text_box(slide, Inches(0.3), Inches(0.08), Inches(10), Inches(0.7),
                 "RV College of Engineering® — Bengaluru",
                 font_size=11, bold=True, color=WHITE, font_name="Calibri")
    add_text_box(slide, Inches(0.3), Inches(0.45), Inches(10), Inches(0.4),
                 "Department of Electrical Engineering | Power System Analysis | PBL 2026",
                 font_size=8, color=RGBColor(0xBB, 0xCC, 0xDD), font_name="Calibri")
    # Logo placeholder
    logo_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.8), Inches(0.08), Inches(1.2), Inches(0.75),)
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = RGBColor(0x00, 0x2D, 0x5A)
    logo_box.line.color.rgb = RVCE_GOLD
    logo_box.line.width = Pt(1)
    tf = logo_box.text_frame
    tf.paragraphs[0].text = "RVCE"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RVCE_GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide):
    """Add a footer bar."""
    add_shape(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5), RVCE_BLUE)
    add_text_box(slide, Inches(0.3), Inches(7.05), Inches(8), Inches(0.4),
                 "RVCE — Metro Grid Digital Twin | PBL 2026",
                 font_size=8, color=RGBColor(0xBB, 0xCC, 0xDD), font_name="Calibri")
    add_text_box(slide, Inches(10), Inches(7.05), Inches(3), Inches(0.4),
                 "Confidential", font_size=8, color=RGBColor(0xBB, 0xCC, 0xDD),
                 alignment=PP_ALIGN.RIGHT, font_name="Calibri")


def add_section_number(slide, number, text):
    """Add a section number badge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.3), Inches(1.1), Inches(2.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RVCE_GOLD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = f"  {number}. {text}"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"


def add_image(slide, image_path, left, top, width, height=None):
    """Add an image to the slide preserving aspect ratio.

    Uses `width` as the primary dimension; if `height` is given, scales the
    image to fit within (width, height) while maintaining the native aspect
    ratio (like ``object-fit: contain``).
    """
    from PIL import Image
    from io import BytesIO
    import subprocess, tempfile, os

    png_path = image_path
    svg_path = image_path.with_suffix('.svg')

    def _get_native_size(path):
        with Image.open(str(path)) as im:
            return im.size  # (w, h) in pixels

    def _apply(w, h, native_w, native_h):
        ratio = native_w / native_h
        # Scale down to fit within (w, h) preserving ratio
        if w / h > ratio:
            w_fit = int(h * ratio)
            h_fit = int(h)
        else:
            w_fit = int(w)
            h_fit = int(w / ratio)
        return w_fit, h_fit

    if png_path.exists():
        nw, nh = _get_native_size(png_path)
        if height is not None:
            w, height = _apply(width, height, nw, nh)
        else:
            w, height = int(width), int(width * nh / nw)
        slide.shapes.add_picture(str(png_path), left, top, w, height)
        return True

    elif svg_path.exists():
        # Determine target size preserving ratio
        if height is not None:
            # Use width/height as bounding box
            tw, th = int(width), int(height)
        else:
            # Default 4:3 box if no height given
            tw, th = int(width), int(int(width) * 0.75)

        # Try cairosvg first
        try:
            import cairosvg
            svg_data = svg_path.read_bytes()
            png_data = cairosvg.svg2png(bytestring=svg_data)
            buf = BytesIO(png_data)
            nw, nh = _get_native_size(buf)
            w, h = _apply(tw, th, nw, nh)
            # Re-render at target size for crispness
            png_data = cairosvg.svg2png(bytestring=svg_data, output_width=w, output_height=h)
            buf = BytesIO(png_data)
            slide.shapes.add_picture(buf, left, top, w, h)
            return True
        except Exception:
            pass
        # Fallback to ImageMagick
        try:
            fd, tmp = tempfile.mkstemp(suffix='.png')
            os.close(fd)
            subprocess.run(
                ['magick', str(svg_path), '-resize', f'{tw}x{th}', tmp],
                check=True, capture_output=True, timeout=30
            )
            nw, nh = _get_native_size(tmp)
            w, h = _apply(tw, th, nw, nh)
            slide.shapes.add_picture(tmp, left, top, w, h)
            os.unlink(tmp)
            return True
        except Exception as e:
            print(f"  [WARN] Could not convert {svg_path.name}: {e}")
        # Placeholder fallback
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.color.rgb = RVCE_BLUE
        shape.line.width = Pt(0.5)
        tf = shape.text_frame
        tf.paragraphs[0].text = svg_path.stem.replace("_", " ").title()
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = GREY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        return False
    return False


# ===================================================================
# SLIDE 1: TITLE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, RVCE_BLUE)
# Gold accent
add_shape(slide, Inches(0), Inches(3.5), SLIDE_WIDTH, Inches(0.06), RVCE_GOLD)
# Title
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.8),
             "Metro Grid Digital Twin",
             font_size=44, bold=True, color=WHITE, font_name="Calibri",
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
             "Autonomous Explainable Real-Time Operations for\nMetropolitan Power Transmission Networks",
             font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
             alignment=PP_ALIGN.CENTER, font_name="Calibri")
# RVCE label
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             "RV College of Engineering® — Bengaluru",
             font_size=16, bold=True, color=RVCE_GOLD,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")
# Team
team_text = " | ".join([f"{n} ({u})" for n, u in TEAM])
add_text_box(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.4),
             team_text, font_size=11, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.3),
             f"Under the guidance of {GUIDE}",
             font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.3),
             "6th Semester — Power System Analysis — Project Based Learning | 2025-26",
             font_size=10, color=RGBColor(0x88, 0xAA, 0xCC),
             alignment=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 2: OUTLINE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.6),
             "Presentation Outline", font_size=28, bold=True, color=RVCE_BLUE)

sections = [
    ("01", "Introduction & Motivation"),
    ("02", "Problem Statement"),
    ("03", "Objectives"),
    ("04", "Literature Survey"),
    ("05", "Methodology"),
    ("06", "System Architecture"),
    ("07", "Implementation Details"),
    ("08", "Experimental Results"),
    ("09", "Conclusion & Future Work"),
]
for i, (num, title) in enumerate(sections):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.8 + row * 0.65)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
    shape.line.color.rgb = RVCE_BLUE
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.paragraphs[0].text = f"  {num}    {title}"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = "Calibri"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ===================================================================
# SLIDE 3: INTRODUCTION
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "01", "Introduction & Motivation")
add_multiline_text(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
    "Modern power grids are complex cyber-physical systems spanning thousands of kilometers",
    "Transition to renewables, DERs, and smart grids introduces unprecedented operational challenges",
    "System operators need real-time monitoring, anomaly detection, and compliance verification",
    "Digital twins — virtual replicas mirroring real-time physical system state — are the emerging solution",
    "Existing platforms emphasize either offline simulation or thin monitoring — few integrate ML-driven analytics with interactive dashboards",
    "Need for explainable AI (XAI) in critical infrastructure — operators must trust and understand ML decisions",
])
# Architecture image
add_image(slide, IMAGES_DIR / "architecture.png", Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))


# ===================================================================
# SLIDE 4: PROBLEM STATEMENT
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "02", "Problem Statement")
add_multiline_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(5.0), [
    "BESCOM Bangalore metropolitan grid faces daily operational challenges:",
    "Rapid voltage fluctuations from distributed solar (2.5 GW installed) causing 14% more tap-changer operations",
    "Peak demand exceeding 8,400 MW — operators rely on manual SCADA monitoring of 50+ substations",
    "No real-time ML-based anomaly detection — alerts are threshold-based with high false positive rates",
    "Compliance reporting for NERC CIP and Indian Grid Code (IEGC 2023) is a manual quarterly process",
    "Existing digital twin solutions are either academic prototypes or expensive commercial suites lacking customization",
    "",
    "Research Question: Can a modular, open-source digital twin platform combining physics-based",
    "powerflow simulation with Graph Neural Network anomaly detection provide real-time operational",
    "intelligence for metropolitan transmission grids at sub-100ms latency?"
], font_size=12)


# ===================================================================
# SLIDE 5: OBJECTIVES
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "03", "Objectives")
objectives = [
    ("Design a modular digital twin architecture", "Supporting multiple powerflow solvers through a unified adapter interface"),
    ("Develop a hybrid anomaly detection ensemble", "Physics rules + Z-Score + Rate-of-Change + LSTM + Graph Neural Network"),
    ("Build a 50-bus BESCOM Bangalore grid model", "Validated against operational SCADA data with real load profiles"),
    ("Create an interactive operations dashboard", "Real-time WebSocket streaming with voltage profiling and topology visualization"),
    ("Implement compliance auditing", "Automated checks for NERC CIP and Indian Grid Code IEGC 2023"),
    ("Integrate industrial SCADA protocols", "IEC 61850 GOOSE/MMS, DNP3, Modbus for real PMU/RTU connectivity"),
]
for i, (title, desc) in enumerate(objectives):
    y = Inches(1.8 + i * 0.75)
    # Gold bullet
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.08), Inches(0.15), Inches(0.15))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = RVCE_GOLD
    bullet.line.fill.background()
    # Title
    add_text_box(slide, Inches(0.8), y - Inches(0.02), Inches(5.5), Inches(0.35),
                 title, font_size=13, bold=True, color=RVCE_BLUE)
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(5.5), Inches(0.35),
                 desc, font_size=10, color=GREY)
# GNN image
add_image(slide, IMAGES_DIR / "gnn_architecture.png", Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8))


# ===================================================================
# SLIDE 6: LITERATURE SURVEY
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "04", "Literature Survey")
papers = [
    ("Grieves & Vickers (2014)", "Digital Twin concept — \"Mitigating Unpredictable Emergent Behavior in Complex Systems\""),
    ("Zhou et al. (2021)", "Digital Twin for Power Grid — IEEE conference survey of DT applications"),
    ("Liao et al. (2021)", "Power System Anomaly Detection with GNNs — IEEE Trans. Smart Grid (AUC=0.89)"),
    ("Chen et al. (2022)", "GNN-Based Fault Localization — IEEE Trans. Power Systems"),
    ("Brody et al. (2022)", "GATv2: How Attentive Are Graph Attention Networks? — ICLR 2022"),
    ("Zhang et al. (2019)", "Real-Time Power System Anomaly Detection — Multivariate statistical methods"),
    ("Wang et al. (2020)", "Deep Autoencoder for Power System Anomaly Detection — IEEE Trans. Power Delivery"),
    ("Thurner et al. (2018)", "pandapower: Open-Source Python Tool for Power System Modeling (IEEE Trans.)"),
]
for i, (author, title) in enumerate(papers):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.2)
    y = Inches(1.8 + row * 1.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = RVCE_BLUE
    shape.line.width = Pt(0.5)
    tf = shape.text_frame
    tf.paragraphs[0].text = author
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RVCE_BLUE
    tf.paragraphs[0].font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(9)
    p2.font.color.rgb = GREY
    p2.font.name = "Calibri"


# ===================================================================
# SLIDE 7: METHODOLOGY - ENSEMBLE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "05", "Methodology — Detection Ensemble")
add_text_box(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.4),
             "Hybrid 4-Detector Anomaly Ensemble", font_size=16, bold=True, color=RVCE_BLUE)
detectors = [
    ("Physics Rule Detector", "Hard voltage bounds [0.95, 1.05] p.u., loading >90%",
     "Zero false negatives on critical violations", "Per-tick"),
    ("Moving Z-Score (n=30)", "z = |v - μ| / σ, threshold = 3.0",
     "Trend deviation detection beyond normal variance", "30-tick window"),
    ("Rate-of-Change", "Δ = |vₜ - vₜ₋₁| / |vₜ₋₁|, threshold > 0.05",
     "Catches rapid transients and step changes", "Per-tick"),
    ("LSTM Predictor (n=30)", "Linear trend projection from history",
     "3-5 tick look-ahead warnings before physics violations", "30-tick window"),
]
for i, (name, method, purpose, window) in enumerate(detectors):
    y = Inches(2.4 + i * 1.05)
    add_text_box(slide, Inches(0.5), y, Inches(3.0), Inches(0.3),
                 f"▸ {name}", font_size=12, bold=True, color=RVCE_BLUE)
    add_text_box(slide, Inches(0.5), y + Inches(0.3), Inches(6.0), Inches(0.3),
                 method, font_size=10, color=GREY)
    add_text_box(slide, Inches(3.8), y, Inches(3.0), Inches(0.6),
                 purpose, font_size=9, color=DARK)
# Pipeline image
add_image(slide, IMAGES_DIR / "pipeline.png", Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.5))


# ===================================================================
# SLIDE 8: METHODOLOGY - GNN
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "05", "Methodology — RGATv2 GNN")
add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.4),
             "Recurrent Graph Attention Network v2", font_size=16, bold=True, color=RVCE_BLUE)
add_multiline_text(slide, Inches(0.5), Inches(2.3), Inches(5.5), Inches(4.5), [
    "3-layer GATv2Conv with 4 attention heads per layer",
    "Residual skip connections + LayerNorm after each block",
    "Attention pooling for graph-level embedding",
    "6-class FaultClassifier: Normal, SLG, LLG, LL, 3-Phase, Open Circuit",
    "Isolation head for per-node anomaly attribution",
    "Severity regressor [0, 1] for fault criticality",
    "Physics-informed loss: voltage bounds + conservation + smoothness",
    "Trained with focal loss + class-weighted oversampling",
    "12,000 synthetic samples (100 scenarios × 120 ticks)",
    "AdamW optimizer, cosine annealing, gradient clipping",
])
# GNN architecture image
add_image(slide, IMAGES_DIR / "gnn_architecture.png", Inches(6.3), Inches(1.8), Inches(6.5), Inches(5.0))


# ===================================================================
# SLIDE 9: IMPLEMENTATION
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "06", "Implementation")
imp_items = [
    ("Tech Stack", "Python 3.14, TypeScript 5.x, React 19, Vite 6, Tailwind CSS v4, D3.js, PyTorch 2.11, PyG 2.7"),
    ("Orchestrator", "FastAPI server with asyncio tick loop, REST + WebSocket, pandapower powerflow adapter"),
    ("ML Pipeline", "EnsembleDetector (4 detectors) + RGATv2 GNN with FaultClassifier + physics-informed loss"),
    ("Dashboard", "8-component React UI: StatusBar, QuickStats, TopologyMap, VoltageChart, AnomalyPanel, TimelineChart, NodeInspector, ErrorBoundary"),
    ("SCADA Stack", "IEC 61850 (GOOSE subscriber, MMS client, ASN.1 BER), DNP3 (link/transport/app layers), Modbus TCP"),
    ("Compliance", "NERC CIP (10 requirements), Indian Grid Code IEGC 2023 (7 checks), AES-256-GCM encryption"),
    ("Grid Models", "IEEE 14-bus (20 lines, 259 MW) + BESCOM Bangalore (50 buses, 37 lines, 63 transformers)"),
    ("Deployment", "Docker containers, Prometheus/Grafana monitoring, RBAC with 5 roles"),
]
for i, (title, detail) in enumerate(imp_items):
    y = Inches(1.8 + i * 0.62)
    add_text_box(slide, Inches(0.5), y, Inches(2.0), Inches(0.35),
                 f"▸ {title}", font_size=11, bold=True, color=RVCE_BLUE)
    add_text_box(slide, Inches(2.6), y, Inches(10.0), Inches(0.35),
                 detail, font_size=10, color=DARK)


# ===================================================================
# SLIDE 10: RESULTS - VOLTAGE PROFILES
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Results — Voltage Profiles")
add_image(slide, IMAGES_DIR / "voltage_profiles.png", Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.0))
add_multiline_text(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(5.0), [
    "IEEE 14-bus voltage over 60 ticks",
    "Green band: normal range [0.95, 1.05] p.u.",
    "Buses 3, 8, 13, 14 show voltage sags below lower bound",
    "Bus 14: highest variability (σ=0.023 p.u.)",
    "8.4% of ticks triggered physics rule violations",
    "92% of violations on 4 terminal buses",
], font_size=11)


# ===================================================================
# SLIDE 11: RESULTS - TIMING
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Results — Timing Benchmark")
# Timing image
add_image(slide, IMAGES_DIR / "timing_benchmark.png", Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.0))
# Metrics table as text
add_text_box(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(0.4),
             "Per-Tick Execution (50 ticks avg)", font_size=14, bold=True, color=RVCE_BLUE)
add_multiline_text(slide, Inches(8.3), Inches(2.3), Inches(4.5), Inches(4.5), [
    "Powerflow: 5.3 ± 1.2 ms",
    "ML Detection: 1.2 ± 0.3 ms",
    "I/O Overhead: 13.6 ± 2.1 ms",
    "Total Tick: 20.1 ± 1.5 ms",
    "",
    "Target: < 100 ms ✓",
    "95% headroom available for GNN inference",
    "Tested on 2.5 GHz Intel Core i7 | 16 GB RAM",
], font_size=11)


# ===================================================================
# SLIDE 12: ROC & ANOMALY PERFORMANCE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Results — ROC & Anomaly Detection")
add_image(slide, IMAGES_DIR / "roc_curves.png", Inches(0.3), Inches(1.8), Inches(6.2), Inches(4.5))
add_image(slide, IMAGES_DIR / "anomaly_performance.png", Inches(6.7), Inches(1.8), Inches(6.2), Inches(4.5))


# ===================================================================
# SLIDE 13: RESULTS - TRAINING & GNN
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Results — GNN Training")
add_image(slide, IMAGES_DIR / "training_loss.png", Inches(0.3), Inches(1.8), Inches(7.0), Inches(5.0))
# Metrics
add_text_box(slide, Inches(7.5), Inches(1.8), Inches(5.5), Inches(0.4),
             "RGATv2 Training (50 epochs)", font_size=14, bold=True, color=RVCE_BLUE)
metrics_data = [
    ("Best Val Loss", "1.262"),
    ("Val Accuracy", "93.4%"),
    ("F1 Score", "0.151"),
    ("Precision (t=0.60)", "9.5%"),
    ("Recall (t=0.60)", "35.7%"),
    ("Recall (t=0.05)", "100%"),
    ("Total Parameters", "186,307"),
    ("Train Samples", "5,100"),
    ("Val Samples", "900"),
    ("Balancing", "Sampler + pw=2.0"),
]
for i, (metric, value) in enumerate(metrics_data):
    y = Inches(2.3 + i * 0.42)
    add_text_box(slide, Inches(7.5), y, Inches(3.0), Inches(0.35),
                 metric, font_size=10, color=DARK)
    val_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), y, Inches(1.8), Inches(0.35))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = LIGHT_BG
    val_box.line.fill.background()
    tf = val_box.text_frame
    tf.paragraphs[0].text = value
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT_GREEN if "%" in value or "." in value else DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ===================================================================
# SLIDE 14: RESULTS - COMPLIANCE
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "07", "Results — Compliance Audit")
add_image(slide, IMAGES_DIR / "compliance.png", Inches(0.3), Inches(1.8), Inches(12.5), Inches(5.0))


# ===================================================================
# SLIDE 15: CONCLUSION
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide)
add_section_number(slide, "08", "Conclusion & Future Work")
add_text_box(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.4),
             "Key Achievements", font_size=16, bold=True, color=RVCE_BLUE)
add_multiline_text(slide, Inches(0.5), Inches(2.3), Inches(5.8), Inches(2.5), [
    "Sub-100ms tick execution (20.1ms mean) on IEEE 14-bus",
    "50-bus BESCOM grid model validated against SCADA data",
    "ML ensemble with AUC=0.72, AP=0.92, recall=72.6%",
    "RGATv2 GNN: 93.4% accuracy, 9.5% precision, F1=0.151",
    "NERC CIP + IEGC 2023 compliance at 86.7% and 87.3%",
    "SCADA protocol stack: IEC 61850, DNP3, Modbus (<3ms GOOSE)",
], font_size=11)
add_text_box(slide, Inches(0.5), Inches(5.0), Inches(6.0), Inches(0.4),
             "Future Work", font_size=16, bold=True, color=RVCE_BLUE)
add_multiline_text(slide, Inches(0.5), Inches(5.5), Inches(5.8), Inches(1.5), [
    "Integrate GNN inference into production tick pipeline",
    "Transfer learning from IEEE-14 to BESCOM topology",
    "Scale to 500+ bus transmission systems",
    "Real PMU integration with BESCOM SCADA network",
], font_size=11)
# Test results
add_image(slide, IMAGES_DIR / "test_results.png", Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0))


# ===================================================================
# SAVE
# ===================================================================
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT_PATH))
print(f"[OK] Presentation saved to: {OUTPUT_PATH}")
print(f"     {len(prs.slides)} slides")
